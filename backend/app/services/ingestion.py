from io import BytesIO
from pathlib import Path
from uuid import uuid4

from docx import Document as DocxDocument
from fastapi import UploadFile
from pptx import Presentation
from pypdf import PdfReader

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".md"}
UPLOAD_DIR = Path("uploads")


class UnsupportedDocumentTypeError(ValueError):
    pass


class EmptyDocumentError(ValueError):
    pass


def get_file_extension(filename: str) -> str:
    return Path(filename).suffix.lower()


def get_file_type(filename: str) -> str:
    extension = get_file_extension(filename)
    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise UnsupportedDocumentTypeError(
            f"Unsupported file type '{extension}'. Supported: {supported}"
        )
    return extension.removeprefix(".")


async def save_upload_file(file: UploadFile) -> Path:
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    extension = get_file_extension(file.filename or "")
    stored_name = f"{uuid4().hex}{extension}"
    destination = UPLOAD_DIR / stored_name

    content = await file.read()
    destination.write_bytes(content)
    await file.seek(0)

    return destination


async def extract_text_from_upload(file: UploadFile) -> str:
    file_type = get_file_type(file.filename or "")
    content = await file.read()
    await file.seek(0)

    if file_type == "pdf":
        text = extract_pdf_text(content)
    elif file_type == "pptx":
        text = extract_pptx_text(content)
    elif file_type == "docx":
        text = extract_docx_text(content)
    else:
        text = content.decode("utf-8", errors="ignore")

    cleaned = clean_text(text)
    if not cleaned:
        raise EmptyDocumentError("No extractable text found in uploaded document")
    return cleaned


def extract_pdf_text(content: bytes) -> str:
    reader = PdfReader(BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages)


def extract_pptx_text(content: bytes) -> str:
    presentation = Presentation(BytesIO(content))
    slide_text: list[str] = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)

    return "\n\n".join(slide_text)


def extract_docx_text(content: bytes) -> str:
    document = DocxDocument(BytesIO(content))
    return "\n\n".join(paragraph.text for paragraph in document.paragraphs)


def clean_text(text: str) -> str:
    lines = [line.strip() for line in text.replace("\x00", "").splitlines()]
    return "\n".join(line for line in lines if line)


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 150) -> list[str]:
    if chunk_size <= 0:
        raise ValueError("chunk_size must be positive")
    if overlap < 0 or overlap >= chunk_size:
        raise ValueError("overlap must be non-negative and smaller than chunk_size")

    chunks: list[str] = []
    start = 0
    text_length = len(text)

    while start < text_length:
        end = min(start + chunk_size, text_length)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end == text_length:
            break
        start = end - overlap

    return chunks
